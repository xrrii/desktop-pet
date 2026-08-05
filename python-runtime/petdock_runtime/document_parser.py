"""PetDock C4 统一文档解析注册表。

本模块是附件和知识库唯一的格式解析入口。第三方解析库只在 ZIP/XML、PDF
主动内容、图片尺寸和路径安全检查通过后调用；Parser 输出结构块，纯文本由
结构块统一生成，避免预览文本与来源位置漂移。
"""

from __future__ import annotations

import io
import logging
import mimetypes
import re
import stat
import zipfile
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Callable

from defusedxml import ElementTree as SafeET
from PIL import Image, ImageOps

LOGGER = logging.getLogger("petdock.document_parser")

MAX_DOCUMENT_BYTES = 10 * 1024 * 1024
MAX_ARCHIVE_FILES = 2_000
MAX_ARCHIVE_ITEM_BYTES = 25 * 1024 * 1024
MAX_ARCHIVE_TOTAL_BYTES = 100 * 1024 * 1024
MAX_ARCHIVE_RATIO = 1_000
MAX_IMAGE_PIXELS = 40_000_000
MAX_IMAGE_FRAMES = 4
MAX_DERIVED_IMAGE_BYTES = 8 * 1024 * 1024

OFFICE_EXTENSIONS = {".docx", ".xlsx", ".pptx"}
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp", ".tif", ".tiff"}
TEXT_EXTENSIONS = {
    ".txt", ".md", ".markdown", ".json", ".jsonl", ".csv", ".tsv", ".yaml", ".yml",
    ".toml", ".ini", ".conf", ".cfg", ".xml", ".html", ".htm", ".css", ".scss", ".less",
    ".js", ".mjs", ".cjs", ".jsx", ".ts", ".tsx", ".vue", ".svelte", ".py", ".pyi",
    ".java", ".kt", ".kts", ".go", ".rs", ".c", ".h", ".cc", ".cpp", ".hpp", ".cs",
    ".swift", ".dart", ".rb", ".php", ".sql", ".sh", ".bash", ".zsh", ".ps1", ".bat",
    ".cmd", ".properties", ".gradle",
}


@dataclass(frozen=True)
class DocumentProblem:
    """表示不会泄露正文或本地路径的结构化解析问题。"""

    code: str
    message: str
    retryable: bool = False

    def as_dict(self) -> dict[str, object]:
        """转换为跨 Runtime/Renderer 的问题协议。"""
        return {"code": self.code, "message": self.message, "retryable": self.retryable}


@dataclass(frozen=True)
class DocumentLocation:
    """描述结构块在原文件中的可引用位置。"""

    kind: str
    value: str
    page: int | None = None
    heading_path: tuple[str, ...] = ()
    paragraph: int | None = None
    sheet: str | None = None
    cell_range: str | None = None
    slide: int | None = None

    def as_dict(self) -> dict[str, object]:
        """返回稳定的 JSON 位置字段。"""
        return {
            "kind": self.kind,
            "value": self.value,
            "page": self.page,
            "headingPath": list(self.heading_path),
            "paragraph": self.paragraph,
            "sheet": self.sheet,
            "cellRange": self.cell_range,
            "slide": self.slide,
        }


@dataclass(frozen=True)
class StructureBlock:
    """保存一段可注入上下文的正文和其来源位置。"""

    kind: str
    content: str
    location: DocumentLocation
    metadata: dict[str, object] = field(default_factory=dict)

    def as_dict(self) -> dict[str, object]:
        """转换为预览和知识库共享的结构块协议。"""
        return {
            "kind": self.kind,
            "content": self.content,
            "location": self.location.as_dict(),
            "metadata": dict(self.metadata),
        }


@dataclass
class ParsedDocument:
    """统一承载标题、结构块、脱敏元数据、警告和错误。"""

    title: str
    parser_id: str
    blocks: list[StructureBlock]
    metadata: dict[str, object] = field(default_factory=dict)
    warnings: list[DocumentProblem] = field(default_factory=list)
    errors: list[DocumentProblem] = field(default_factory=list)

    @property
    def plain_text(self) -> str:
        """从结构块按原顺序生成规范化纯文本。"""
        return "\n\n".join(block.content for block in self.blocks if block.content.strip()).strip()

    @property
    def ready(self) -> bool:
        """判断文档是否可以进入模型上下文。"""
        return not self.errors

    def as_dict(self) -> dict[str, object]:
        """转换为 SQLite/HTTP 可持久化的脱敏结构。"""
        return {
            "title": self.title,
            "parserId": self.parser_id,
            "plainText": self.plain_text,
            "blocks": [block.as_dict() for block in self.blocks],
            "metadata": dict(self.metadata),
            "warnings": [item.as_dict() for item in self.warnings],
            "errors": [item.as_dict() for item in self.errors],
        }


class DocumentParseError(ValueError):
    """携带稳定错误码的文档解析异常。"""

    def __init__(self, code: str, message: str, *, retryable: bool = False) -> None:
        """初始化解析错误，不保存路径、正文或原始异常消息。"""
        super().__init__(code)
        self.problem = DocumentProblem(code, message, retryable)


def validate_office_zip(path: Path) -> list[str]:
    """在第三方 Office 库读取前检查压缩包大小、路径、重复项和危险对象。"""
    if path.read_bytes()[:8] == b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1":
        raise DocumentParseError("document_encrypted", "Office 文档是加密容器，无法读取。")
    try:
        with zipfile.ZipFile(path) as archive:
            infos = archive.infolist()
            if len(infos) > MAX_ARCHIVE_FILES:
                raise DocumentParseError("document_archive_limit_exceeded", "Office 压缩包文件数量超限。")
            names: set[str] = set()
            total = 0
            warnings: list[str] = []
            for info in infos:
                name = info.filename
                normalized = name.replace("\\", "/")
                if normalized in names:
                    raise DocumentParseError("document_archive_limit_exceeded", "Office 压缩包包含重复文件名。")
                names.add(normalized)
                if normalized.startswith("/") or any(part == ".." for part in normalized.split("/")):
                    raise DocumentParseError("document_archive_unsafe", "Office 压缩包路径不安全。")
                mode = (info.external_attr >> 16) & 0xFFFF
                if stat.S_ISLNK(mode):
                    raise DocumentParseError("document_archive_unsafe", "Office 压缩包包含符号链接。")
                if info.flag_bits & 0x1:
                    raise DocumentParseError("document_encrypted", "Office 文档已加密，无法读取。")
                if info.file_size > MAX_ARCHIVE_ITEM_BYTES:
                    raise DocumentParseError("document_archive_limit_exceeded", "Office 压缩包单项大小超限。")
                total += info.file_size
                if total > MAX_ARCHIVE_TOTAL_BYTES:
                    raise DocumentParseError("document_archive_limit_exceeded", "Office 压缩包解压总大小超限。")
                if info.compress_size and info.file_size / info.compress_size > MAX_ARCHIVE_RATIO:
                    raise DocumentParseError("document_archive_limit_exceeded", "Office 压缩包压缩比超限。")
                lowered = normalized.casefold()
                if any(token in lowered for token in ("vba", "activex", "embeddings/", "oleobject")):
                    warnings.append("document_active_content_ignored")
                if lowered.endswith("externalrelationships.xml") or "externalLinks/".casefold() in lowered:
                    warnings.append("document_external_link_ignored")
            # 所有 XML 必须在 python-docx/openpyxl/python-pptx 前完成 DTD/ENTITY 扫描。
            for info in infos:
                lowered = info.filename.casefold()
                if not (lowered.endswith(".xml") or lowered.endswith(".rels")):
                    continue
                xml_root = _safe_xml(archive.read(info))
                if lowered.endswith(".rels"):
                    for relationship in xml_root.iter():
                        target_mode = str(relationship.attrib.get("TargetMode", "")).casefold()
                        target = str(relationship.attrib.get("Target", "")).casefold()
                        rel_type = str(relationship.attrib.get("Type", "")).casefold()
                        if target_mode == "external" or target.startswith(("http:", "https:", "file:")):
                            warnings.append("document_external_link_ignored")
                        if any(token in rel_type for token in ("attachedtemplate", "oleobject", "package")):
                            warnings.append("document_active_content_ignored")
            return sorted(set(warnings))
    except zipfile.BadZipFile as error:
        raise DocumentParseError("document_corrupt", "Office 文档压缩结构损坏。") from error


def _safe_xml(raw: bytes) -> Any:
    """拒绝 DTD/ENTITY 后使用 defusedxml 解析 XML。"""
    if re.search(rb"<!DOCTYPE|<!ENTITY", raw, re.IGNORECASE):
        raise DocumentParseError("document_xml_unsafe", "文档 XML 包含禁止的 DTD 或 ENTITY。")
    try:
        return SafeET.fromstring(raw)
    except Exception as error:
        raise DocumentParseError("document_corrupt", "文档 XML 结构损坏。") from error


def _problem_for_warnings(values: list[str]) -> list[DocumentProblem]:
    """把安全扫描的内部标记转换为用户可理解的结构化告警。"""
    return [DocumentProblem(code, "文档中的主动内容或外部链接已忽略。") for code in sorted(set(values))]


def _text(raw: bytes) -> str:
    """严格解码文本附件，避免系统编码猜测造成乱码。"""
    try:
        return raw.decode("utf-8-sig")
    except UnicodeDecodeError as error:
        raise DocumentParseError("attachment_decode_failed", "文本附件不是有效的 UTF-8。") from error


def _image_document(path: Path, *, vision_status: str, defer_vision: bool = False) -> ParsedDocument:
    """安全解码图片并返回尺寸、格式和脱敏元数据；不执行 OCR。"""
    try:
        with Image.open(path) as image:
            width, height = image.size
            frames = int(getattr(image, "n_frames", 1))
            image.verify()
            if width * height > MAX_IMAGE_PIXELS or frames > MAX_IMAGE_FRAMES:
                raise DocumentParseError("image_too_large", "图片像素或帧数超限。")
            format_name = (image.format or "").upper()
            metadata = {"format": format_name, "width": width, "height": height, "frames": frames}
    except DocumentParseError:
        raise
    except Exception as error:
        raise DocumentParseError("image_decode_failed", "图片无法安全解码。") from error
    error_code = {
        "unconfigured": "vision_not_configured",
        "untested": "vision_capability_untested",
        "unsupported": "vision_model_unsupported",
        "unavailable": "vision_provider_unavailable",
        "invalid-credentials": "vision_invalid_credentials",
    }.get(vision_status)
    errors = [] if defer_vision or vision_status == "supported" else [DocumentProblem(error_code or "vision_capability_untested", "视觉能力尚不可用于图片附件。")]
    blocks = [StructureBlock("image_metadata", f"图片尺寸：{width} x {height}，格式：{format_name}。", DocumentLocation("image", "metadata"), metadata)]
    return ParsedDocument(path.stem, "image-metadata-v1", blocks, metadata, errors=errors)


def derive_safe_image(path: Path, destination: Path) -> dict[str, object]:
    """移除 EXIF/GPS 并重新编码为无元数据 PNG，返回脱敏派生图信息。"""
    try:
        with Image.open(path) as source:
            if source.size[0] * source.size[1] > MAX_IMAGE_PIXELS or int(getattr(source, "n_frames", 1)) > MAX_IMAGE_FRAMES:
                raise DocumentParseError("image_too_large", "图片像素或帧数超限。")
            image = ImageOps.exif_transpose(source).convert("RGBA")
            destination.parent.mkdir(parents=True, exist_ok=True)
            image.save(destination, format="PNG", optimize=True)
        if destination.stat().st_size > MAX_DERIVED_IMAGE_BYTES:
            destination.unlink(missing_ok=True)
            raise DocumentParseError("image_too_large", "安全派生图大小超限。")
        LOGGER.info("图片安全派生完成 extension=%s bytes=%s", path.suffix.lower(), destination.stat().st_size)
        return {"format": "PNG", "sizeBytes": destination.stat().st_size}
    except DocumentParseError:
        raise
    except Exception as error:
        raise DocumentParseError("image_decode_failed", "图片安全派生失败。") from error


def _pdf_document(path: Path) -> ParsedDocument:
    """读取 PDF 文本层和页码，并拒绝脚本、自动动作及嵌入对象。"""
    from pypdf import PdfReader

    raw = path.read_bytes()
    dangerous = re.search(rb"/(?:JavaScript|JS|AA|OpenAction|Launch|EmbeddedFile|RichMedia|SubmitForm)\b", raw)
    if dangerous:
        raise DocumentParseError("document_active_content", "PDF 包含禁止的脚本或嵌入主动内容。")
    try:
        reader = PdfReader(io.BytesIO(raw), strict=False)
    except Exception as error:
        raise DocumentParseError("document_corrupt", "PDF 结构损坏。") from error
    if reader.is_encrypted:
        raise DocumentParseError("document_encrypted", "PDF 已加密，无法读取。")
    title = str((reader.metadata or {}).get("/Title") or path.stem)[:255]
    blocks: list[StructureBlock] = []
    for page_number, page in enumerate(reader.pages, start=1):
        try:
            content = (page.extract_text() or "").strip()
        except Exception as error:
            LOGGER.warning("PDF 页面解析失败 page=%s error=%s", page_number, error.__class__.__name__)
            content = ""
        if content:
            blocks.append(StructureBlock("paragraph", content, DocumentLocation("pdf_page", f"page-{page_number}", page=page_number)))
    if reader.pages and not blocks:
        return ParsedDocument(title, "pdf-text-v1", [], {"pages": len(reader.pages)}, errors=[DocumentProblem("document_ocr_required", "PDF 没有可读取的文本层，需要 OCR。")])
    return ParsedDocument(title, "pdf-text-v1", blocks, {"pages": len(reader.pages)})


def _docx_document(path: Path) -> ParsedDocument:
    """读取 DOCX 标题、段落、列表和表格，不加载宏或嵌入对象。"""
    warnings = validate_office_zip(path)
    with zipfile.ZipFile(path) as archive:
        root = _safe_xml(archive.read("word/document.xml"))
    ns = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
    blocks: list[StructureBlock] = []
    heading: list[str] = []
    paragraph_number = 0
    table_number = 0
    for node in root.findall(".//w:body/*", ns):
        tag = node.tag.rsplit("}", 1)[-1]
        if tag == "p":
            paragraph_number += 1
            text = "".join(item.text or "" for item in node.findall(".//w:t", ns)).strip()
            if not text:
                continue
            style = node.find("./w:pPr/w:pStyle", ns)
            style_id = style.get(f"{{{ns['w']}}}val", "") if style is not None else ""
            is_heading = style_id.lower().startswith("heading") or style_id.lower().startswith("title")
            if is_heading:
                level_match = re.search(r"(\d+)$", style_id)
                level = int(level_match.group(1)) if level_match else 1
                heading = heading[: max(0, level - 1)] + [text]
                kind = "heading"
            else:
                num = node.find("./w:pPr/w:numPr", ns)
                kind = "list_item" if num is not None or style_id.lower().startswith("list") else "paragraph"
            blocks.append(StructureBlock(kind, text, DocumentLocation("docx_paragraph", f"paragraph-{paragraph_number}", paragraph=paragraph_number, heading_path=tuple(heading))))
        elif tag == "tbl":
            table_number += 1
            rows: list[str] = []
            for row in node.findall("./w:tr", ns):
                rows.append(" | ".join("".join(item.text or "" for item in cell.findall(".//w:t", ns)).strip() for cell in row.findall("./w:tc", ns)))
            content = "\n".join(row for row in rows if row.strip())
            if content:
                blocks.append(StructureBlock("table", content, DocumentLocation("docx_table", f"table-{table_number}", heading_path=tuple(heading))))
    title = heading[0] if heading else path.stem
    return ParsedDocument(title, "docx-ooxml-v1", blocks, {"headingPath": heading}, warnings=_problem_for_warnings(warnings))


def _xlsx_document(path: Path) -> ParsedDocument:
    """读取 XLSX 工作表有效区域、单元格位置和公式文本，不执行公式。"""
    import openpyxl

    warnings = validate_office_zip(path)
    try:
        workbook = openpyxl.load_workbook(path, read_only=True, data_only=False, keep_links=False)
    except RuntimeError as error:
        if "encrypt" in str(error).lower():
            raise DocumentParseError("document_encrypted", "XLSX 已加密，无法读取。") from error
        raise DocumentParseError("document_corrupt", "XLSX 结构损坏。") from error
    blocks: list[StructureBlock] = []
    total_cells = 0
    try:
        for sheet in workbook.worksheets:
            rows: list[str] = []
            sheet_cells = 0
            for row in sheet.iter_rows():
                cells: list[str] = []
                for cell in row:
                    value = cell.value
                    if value is None:
                        continue
                    total_cells += 1
                    sheet_cells += 1
                    if total_cells > 100_000:
                        raise DocumentParseError("document_archive_limit_exceeded", "XLSX 有效单元格数量超限。")
                    cells.append(f"{cell.coordinate}={value}")
                if cells:
                    rows.append(" | ".join(cells))
            if rows:
                blocks.append(StructureBlock("worksheet", "\n".join(rows), DocumentLocation("xlsx_sheet", sheet.title, sheet=sheet.title, cell_range=sheet.calculate_dimension()), {"cellCount": sheet_cells}))
    finally:
        workbook.close()
    return ParsedDocument(path.stem, "xlsx-openpyxl-v1", blocks, {"worksheets": len(blocks), "cellCount": total_cells}, warnings=_problem_for_warnings(warnings))


def _pptx_document(path: Path) -> ParsedDocument:
    """读取 PPTX 每页标题、正文和备注，不执行动画、宏或嵌入对象。"""
    warnings = validate_office_zip(path)
    from pptx import Presentation

    try:
        presentation = Presentation(path)
    except Exception as error:
        raise DocumentParseError("document_corrupt", "PPTX 结构损坏。") from error
    blocks: list[StructureBlock] = []
    for index, slide in enumerate(presentation.slides, start=1):
        title = slide.shapes.title.text.strip() if slide.shapes.title is not None else ""
        body: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape is not slide.shapes.title:
                body.append(shape.text.strip())
        text = "\n".join(item for item in ([title] + body) if item)
        if text:
            blocks.append(StructureBlock("slide", text, DocumentLocation("pptx_slide", f"slide-{index}", slide=index), {"title": title}))
        try:
            notes = slide.notes_slide.notes_text_frame.text.strip()
        except Exception:
            notes = ""
        if notes:
            blocks.append(StructureBlock("notes", notes, DocumentLocation("pptx_slide_notes", f"slide-{index}-notes", slide=index)))
    title = blocks[0].metadata.get("title") if blocks and blocks[0].kind == "slide" else path.stem
    return ParsedDocument(str(title or path.stem), "pptx-python-v1", blocks, {"slides": len(presentation.slides)}, warnings=_problem_for_warnings(warnings))


class DocumentParserRegistry:
    """按文件签名和扩展名选择唯一的 C4 Parser。"""

    def __init__(self, *, vision_enabled: bool = False, vision_status: str | None = None) -> None:
        """创建 Registry；图片只有在 Vision 主动探测成功后才允许进入上下文。"""
        self.vision_status = vision_status or ("supported" if vision_enabled else "untested")

    @property
    def vision_enabled(self) -> bool:
        """兼容旧调用方的布尔能力查询。"""
        return self.vision_status == "supported"

    @vision_enabled.setter
    def vision_enabled(self, value: bool) -> None:
        """兼容旧调用方设置；新代码应直接传递完整状态。"""
        self.vision_status = "supported" if value else "untested"

    def capabilities(self) -> list[dict[str, object]]:
        """返回附件和知识库共用的 Parser 能力声明。"""
        return [
            {"parserId": "pdf-text-v1", "extensions": [".pdf"], "maxInputBytes": MAX_DOCUMENT_BYTES},
            {"parserId": "docx-ooxml-v1", "extensions": [".docx"], "maxInputBytes": MAX_DOCUMENT_BYTES},
            {"parserId": "xlsx-openpyxl-v1", "extensions": [".xlsx"], "maxInputBytes": MAX_DOCUMENT_BYTES},
            {"parserId": "pptx-python-v1", "extensions": [".pptx"], "maxInputBytes": MAX_DOCUMENT_BYTES},
            {"parserId": "image-metadata-v1", "extensions": sorted(IMAGE_EXTENSIONS), "visionRequired": True, "visionEnabled": self.vision_enabled},
            {"parserId": "utf8-text-v1", "extensions": sorted(TEXT_EXTENSIONS), "maxInputBytes": MAX_DOCUMENT_BYTES},
        ]

    def parse(
        self,
        path: str | Path,
        *,
        name: str | None = None,
        mime: str | None = None,
        vision_status: str | None = None,
        defer_vision: bool = False,
    ) -> ParsedDocument:
        """安全解析单个文件；调用级视觉状态用于保持知识库图片索引默认关闭。"""
        source = Path(path)
        if not source.is_file() or source.is_symlink():
            raise DocumentParseError("attachment_not_regular_file", "附件不是普通文件。")
        size = source.stat().st_size
        if size < 1 or size > MAX_DOCUMENT_BYTES:
            raise DocumentParseError("attachment_too_large", "附件大小超出限制。")
        extension = Path(name or source.name).suffix.lower()
        if extension == ".pdf":
            parsed = _pdf_document(source)
        elif extension == ".docx":
            parsed = _docx_document(source)
        elif extension == ".xlsx":
            parsed = _xlsx_document(source)
        elif extension == ".pptx":
            parsed = _pptx_document(source)
        elif extension in IMAGE_EXTENSIONS or (mime or "").startswith("image/"):
            parsed = _image_document(
                source,
                vision_status=vision_status or self.vision_status,
                defer_vision=defer_vision,
            )
        elif extension in TEXT_EXTENSIONS:
            content = _text(source.read_bytes())
            parsed = ParsedDocument(source.stem, "utf8-text-v1", [StructureBlock("text", content, DocumentLocation("text", "document"))], {"characters": len(content)})
        else:
            raise DocumentParseError("attachment_type_unsupported", "附件类型不受支持。")
        LOGGER.info("文档解析完成 extension=%s bytes=%s parser=%s errors=%s", extension, size, parsed.parser_id, len(parsed.errors))
        return parsed
