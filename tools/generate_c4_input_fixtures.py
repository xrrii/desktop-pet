"""为独立 Runtime smoke 生成 C4 正常输入样本。"""

from __future__ import annotations

import json
import sys
from pathlib import Path

import openpyxl
from docx import Document
from PIL import Image
from pptx import Presentation
from pptx.util import Inches


def _write_pdf(path: Path) -> None:
    """生成带文本层的最小 PDF，不引入 C4 禁止的输出依赖。"""
    stream = b"BT /F1 12 Tf 72 720 Td (Packaged PDF C4) Tj ET"
    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] /Resources << /Font << /F1 4 0 R >> >> /Contents 5 0 R >>",
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
        b"<< /Length " + str(len(stream)).encode("ascii") + b" >>\nstream\n" + stream + b"\nendstream",
    ]
    output = bytearray(b"%PDF-1.4\n")
    offsets: list[int] = []
    for index, value in enumerate(objects, start=1):
        offsets.append(len(output))
        output.extend(f"{index} 0 obj\n".encode("ascii") + value + b"\nendobj\n")
    xref = len(output)
    output.extend(f"xref\n0 {len(objects) + 1}\n0000000000 65535 f \n".encode("ascii"))
    output.extend(b"".join(f"{offset:010d} 00000 n \n".encode("ascii") for offset in offsets))
    output.extend(f"trailer\n<< /Size {len(objects) + 1} /Root 1 0 R >>\nstartxref\n{xref}\n%%EOF\n".encode("ascii"))
    path.write_bytes(output)


def _write_docx(path: Path) -> None:
    """生成包含标题、列表和表格的 DOCX。"""
    document = Document()
    document.add_heading("打包 DOCX", level=1)
    document.add_paragraph("只读输入", style="List Bullet")
    table = document.add_table(rows=1, cols=2)
    table.cell(0, 0).text = "键"
    table.cell(0, 1).text = "值"
    document.save(path)


def _write_xlsx(path: Path) -> None:
    """生成包含公式文本的 XLSX。"""
    workbook = openpyxl.Workbook()
    sheet = workbook.active
    sheet.title = "打包数据"
    sheet["A1"] = 2
    sheet["B1"] = "=A1*2"
    workbook.save(path)
    workbook.close()


def _write_pptx(path: Path) -> None:
    """生成包含标题和正文的 PPTX。"""
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "打包 PPTX"
    box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(1))
    box.text_frame.text = "幻灯片正文"
    presentation.save(path)


def generate(root: Path) -> list[dict[str, object]]:
    """在受控附件目录生成样本并返回登记协议。"""
    samples = [
        ("1" * 32, "sample.pdf", _write_pdf),
        ("2" * 32, "sample.docx", _write_docx),
        ("3" * 32, "sample.xlsx", _write_xlsx),
        ("4" * 32, "sample.pptx", _write_pptx),
        ("5" * 32, "sample.png", lambda path: Image.new("RGB", (32, 32), "white").save(path)),
    ]
    registrations: list[dict[str, object]] = []
    for attachment_id, name, writer in samples:
        directory = root / attachment_id
        directory.mkdir(parents=True, exist_ok=True)
        path = directory / name
        writer(path)
        registrations.append({
            "id": attachment_id,
            "name": name,
            "relativePath": f"{attachment_id}/{name}",
            "sizeBytes": path.stat().st_size,
        })
    return registrations


def main() -> None:
    """读取目标目录并向 stdout 输出 UTF-8 JSON 登记列表。"""
    if len(sys.argv) != 2:
        raise SystemExit("用法：generate_c4_input_fixtures.py <attachment-root>")
    print(json.dumps(generate(Path(sys.argv[1])), ensure_ascii=False))


if __name__ == "__main__":
    main()
